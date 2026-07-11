# Copyright 2026 Darwin-Agent
# SPDX-License-Identifier: MIT
from __future__ import annotations

import os

from ..base import tool

_SCHEMA = {
    "type": "object",
    "properties": {
        "file_path": {
            "type": "string",
            "description": "The absolute path to the file to read",
        },
        "limit": {
            "type": "integer",
            "description": "Max number of lines to read (default 2000)",
        },
        "offset": {
            "type": "integer",
            "description": "Line number to start reading from (0-indexed)",
        },
        "pages": {
            "type": "string",
            "description": (
                "Page range for PDF files (e.g. '1-5', '3', '10-20'). "
                "Only applicable to .pdf files. Maximum 20 pages per request. "
                "Requires: pip install 'harnessx'"
            ),
        },
    },
    "required": ["file_path"],
}

# File extensions handled by specialized readers (case-insensitive)
_OFFICE_EXTENSIONS = {".docx", ".xlsx", ".xls", ".csv", ".pptx"}

_MAX_PDF_PAGES = 20


def _continuation_notice(start: int, end: int, total: int) -> str:
    return f"\n[Showing lines {start + 1}–{end} of {total} total. Use offset={end} to read the next section.]"


@tool(
    name="Read",
    description=(
        "Read a file from the local filesystem. Returns file contents with line numbers. "
        "Supports text files, PDFs, DOCX, XLSX/CSV spreadsheets, and PPTX presentations."
    ),
    input_schema=_SCHEMA,
)
async def read_tool(file_path: str, limit: int = 2000, offset: int = 0, pages: str | None = None) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return _read_pdf(file_path, pages)
    if ext in _OFFICE_EXTENSIONS:
        return _read_office(file_path, ext)
    from ...sandbox.base import get_current_sandbox

    sandbox = get_current_sandbox()
    resolved = sandbox.resolve(file_path) if sandbox is not None else file_path
    try:
        if sandbox is not None:
            content = await sandbox.read_file(resolved)
            lines = content.splitlines(keepends=True)
        else:
            with open(resolved, "r", encoding="utf-8", errors="replace") as f:
                lines = f.readlines()
        start = offset
        end = offset + limit if limit > 0 else len(lines)
        selected = lines[start:end]
        result = "".join(f"{start + i + 1}\t{line}" for i, line in enumerate(selected))
        if limit > 0 and end < len(lines):
            result += _continuation_notice(start, end, len(lines))
        return result
    except FileNotFoundError:
        return f"Error: File not found: {file_path}"
    except Exception as e:
        return f"Error: {e}"


def _parse_page_range(pages: str, total: int) -> list[int]:
    """Parse a page range string like '1-5' or '3' into a list of 0-indexed page numbers."""
    pages = pages.strip()
    if "-" in pages:
        parts = pages.split("-", 1)
        start = max(1, int(parts[0].strip()))
        end = min(total, int(parts[1].strip()))
    else:
        start = end = int(pages.strip())
    # Clamp to _MAX_PDF_PAGES
    if end - start + 1 > _MAX_PDF_PAGES:
        end = start + _MAX_PDF_PAGES - 1
    return list(range(start - 1, end))  # 0-indexed


def _read_pdf(file_path: str, pages: str | None) -> str:
    try:
        import pdfplumber  # type: ignore
    except ImportError:
        return "Error: PDF reading requires: pip install 'harnessx'"

    try:
        with pdfplumber.open(file_path) as pdf:
            total = len(pdf.pages)
            if pages:
                page_indices = _parse_page_range(pages, total)
            else:
                if total > _MAX_PDF_PAGES:
                    return (
                        f"Error: PDF has {total} pages. "
                        f"Provide the pages parameter to read specific pages "
                        f"(e.g. pages='1-{_MAX_PDF_PAGES}'). Maximum {_MAX_PDF_PAGES} pages per request."
                    )
                page_indices = list(range(total))

            parts: list[str] = []
            for i in page_indices:
                if i >= total:
                    continue
                text = pdf.pages[i].extract_text() or ""
                parts.append(f"--- Page {i + 1} ---\n{text}")
            return "\n\n".join(parts)
    except FileNotFoundError:
        return f"Error: File not found: {file_path}"
    except Exception as e:
        return f"Error reading PDF: {e}"


def _read_office(file_path: str, ext: str) -> str:
    """Read DOCX, XLSX/XLS/CSV, or PPTX files."""
    try:
        if ext == ".docx":
            return _read_docx(file_path)
        elif ext in (".xlsx", ".xls"):
            return _read_spreadsheet(file_path, ext)
        elif ext == ".csv":
            return _read_csv(file_path)
        elif ext == ".pptx":
            return _read_pptx(file_path)
        else:
            return f"Error: Unsupported file type: {ext}"
    except FileNotFoundError:
        return f"Error: File not found: {file_path}"
    except ImportError as e:
        return f"Error: Missing dependency — {e}. pip install python-docx openpyxl python-pptx"
    except Exception as e:
        return f"Error reading {ext} file: {e}"


def _read_docx(file_path: str) -> str:
    """Read a DOCX file and return its text content."""
    import docx

    doc = docx.Document(file_path)
    parts: list[str] = []
    for i, para in enumerate(doc.paragraphs):
        text = para.text.strip()
        if text:
            parts.append(f"{i + 1}\t{text}")

    # Also read tables
    for t_idx, table in enumerate(doc.tables):
        parts.append(f"\n--- Table {t_idx + 1} ---")
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            parts.append("\t".join(cells))

    return "\n".join(parts) if parts else "(empty document)"


def _read_spreadsheet(file_path: str, ext: str) -> str:
    """Read XLSX/XLS file and return all sheets as text tables."""
    import openpyxl

    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    parts: list[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"--- Sheet: {sheet_name} ---")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            parts.append("\t".join(cells))
    wb.close()
    return "\n".join(parts) if parts else "(empty workbook)"


def _read_csv(file_path: str) -> str:
    """Read a CSV file and return its content."""
    import csv

    parts: list[str] = []
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f)
        for i, row in enumerate(reader):
            parts.append(f"{i + 1}\t" + "\t".join(row))
            if i >= 5000:  # safety limit
                parts.append(f"\n[... truncated at {i + 1} rows]")
                break
    return "\n".join(parts) if parts else "(empty file)"


def _read_pptx(file_path: str) -> str:
    """Read a PPTX file and return slide text."""
    from pptx import Presentation

    prs = Presentation(file_path)
    parts: list[str] = []
    for i, slide in enumerate(prs.slides):
        parts.append(f"--- Slide {i + 1} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)
    return "\n".join(parts) if parts else "(empty presentation)"
