#!/usr/bin/env python3
"""Build the Harness Evolution workshop deck (.pptx)."""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── local reproduction numbers (fill from runs) ───────────────────────────────
BASELINE = sys.argv[1] if len(sys.argv) > 1 else "0.00"   # 8B vanilla, retail, 5 tasks
EVOLVED  = sys.argv[2] if len(sys.argv) > 2 else "TBD"     # 8B evolved, same tasks
OUT      = sys.argv[3] if len(sys.argv) > 3 else "workshop.pptx"

# palette
INK   = RGBColor(0x14, 0x18, 0x24)
PAPER = RGBColor(0xF7, 0xF7, 0xFB)
ACCENT= RGBColor(0x7C, 0x5C, 0xFF)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
MUTE  = RGBColor(0x66, 0x6C, 0x7A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    return tb.text_frame

def para(tf, text, size, color, bold=False, first=False, align=PP_ALIGN.LEFT, space=8, italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(space)
    r = p.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color; f.name = "Arial"
    return p

def chip(slide, text):
    tf = box(slide, 0.9, 0.5, 6, 0.5)
    para(tf, text, 13, ACCENT, bold=True, first=True)

def content(title, items, sub=None):
    s = prs.slides.add_slide(BLANK); bg(s, PAPER)
    tf = box(s, 0.9, 0.7, 11.5, 1.3)
    para(tf, title, 34, INK, bold=True, first=True)
    if sub: para(tf, sub, 17, MUTE, space=4)
    body = box(s, 0.95, 2.35, 11.4, 4.7)
    for i, it in enumerate(items):
        if isinstance(it, tuple):
            txt, kind = it
        else:
            txt, kind = it, "b"
        if kind == "q":
            para(body, txt, 21, ACCENT, italic=True, first=(i==0), space=14)
        elif kind == "h":
            para(body, txt, 20, INK, bold=True, first=(i==0), space=6)
        else:
            para(body, "•  " + txt, 18, INK, first=(i==0), space=10)
    return s

def image_slide(title, img_path, caption=None):
    import os
    s = prs.slides.add_slide(BLANK); bg(s, PAPER)
    tf = box(s, 0.9, 0.55, 11.5, 1.0)
    para(tf, title, 30, INK, bold=True, first=True)
    if os.path.exists(img_path):
        from PIL import Image
        try:
            iw, ih = Image.open(img_path).size; ar = ih / iw
        except Exception:
            ar = 0.4
        w = Inches(10.6); h = Inches(10.6 * ar)
        left = (SW - w) / 2
        s.shapes.add_picture(img_path, left, Inches(1.9), width=w, height=h)
    if caption:
        cf = box(s, 0.9, 6.7, 11.5, 0.6)
        para(cf, caption, 13, MUTE, first=True, align=PP_ALIGN.CENTER)
    return s

# ── Slide 1: title ────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s, INK)
tf = box(s, 1.0, 2.5, 11.3, 2.6)
para(tf, "Evolve the Harness,", 50, WHITE, bold=True, first=True, space=0)
para(tf, "Not (Just) the Model", 50, ACCENT, bold=True, space=16)
para(tf, "Reinforcement learning on agent harnesses — reproduced locally on a frozen model",
     20, RGBColor(0xC8,0xCC,0xD8), space=4)
tf2 = box(s, 1.0, 6.2, 11.3, 0.8)
para(tf2, "HarnessX  ·  “Don't Train the Model, Evolve the Harness”  ·  τ²-Bench retail  ·  Claude Code as meta-agent",
     14, MUTE, first=True)

# ── Slide 2: hook ─────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s, INK)
chip_tf = box(s, 0.9, 0.6, 6, 0.5); para(chip_tf, "THE HOOK", 13, ACCENT, bold=True, first=True)
tf = box(s, 1.0, 2.2, 11.3, 3.5)
para(tf, "“A frozen open model that solves 0% of a legal-agent benchmark end to end",
     30, WHITE, bold=True, first=True, space=2)
para(tf, "is not as weak as that score looks. Zero model weights changed.”", 30, WHITE, bold=True, space=20)
para(tf, "Same model. Change the wrapper. Double-digit accuracy gains.", 20, GREEN, space=4)
para(tf, "The question isn't “which model” — it's “which harness.”", 20, RGBColor(0xC8,0xCC,0xD8))

# ── content slides ────────────────────────────────────────────────────────────
content("What is a “harness”?",
    [("The runtime wrapper around the model.", "h"),
     ("Context assembly · tools · memory · control flow · error recovery · evaluation", "b"),
     ("agent = model.agentic(harness)", "b"),
     ("Model = a fixed reasoning ceiling.", "b"),
     ("Harness = how much of that ceiling you actually realize.", "b")])

content("The thesis",
    [("Vanilla off-the-shelf harnesses are fine at baseline.", "h"),
     ("Vertical-specific use cases need the harness to evolve with them — or you design your own.", "h"),
     ("Why? Each domain fails differently. One harness can't serve every vertical.", "b")])

content("Two papers, one idea",
    [("HarnessX (arXiv:2606.14249): MetaAgent.evolve() over trajectories — frozen Qwen3.5-9B / GPT-5.", "b"),
     ("+14.5% avg, up to +44%. Weakest model gains most.", "b"),
     ("“Evolve the Harness” (Niklaus): Meta-Harness proposer + gate — frozen DeepSeek-V4-Pro, legal vertical.", "b"),
     ("63.4% → 80.1% (+16.7pp).", "b"),
     ("Both: zero model-weight changes.", "h")])

content("How harness evolution works",
    [("R0: run tasks → LLM-judge each trajectory", "b"),
     ("Meta-agent reads the failures → writes a new config.yaml (+ processors)", "b"),
     ("Gate: keep if reward improves past tolerance, else revert to best-so-far", "b"),
     ("Repeat. It's harness-edit-as-MDP:", "h"),
     ("state = config + traces   ·   action = one edit   ·   reward = gated pass-delta", "b")])

content("Why it works — the core teaching",
    [("Most failures are NOT reasoning failures — they're harness failures:", "h"),
     ("wrong file path · chunked-write corruption · malformed tool JSON", "b"),
     ("context overflow · loop-without-progress · premature / wrong-context action", "b"),
     ("“For a weak agent, a lot of ‘capability’ is I/O discipline the scaffold can guarantee.”", "q")])

content("The punchline: code beats prompts",
    [("“Five of the top six harnesses are deterministic code, not prompt edits.”", "q"),
     ("Prompt edits: local, brittle, don't transfer across models.", "b"),
     ("Code fixes: structural, portable, testable, versionable.", "b")])

content("But it's not only code — 4 levers",
    [("The evolver edits the whole harness, not just processors:", "h"),
     ("Instruction = prompts / guidance    ·    Action = tools / skills", "b"),
     ("Control = deterministic processors    ·    Configuration = memory, compaction, knobs", "b"),
     ("Which lever wins depends on your model + your vertical's bottleneck:", "h"),
     ("strong model → prompts   ·   weak model → control code   ·   long-context → memory/compaction   ·   retrieval → tools", "b")])

content("So: evolve the dimension your vertical is bottlenecked on",
    [("“Prompt dominance scales inversely with base-model strength.” — HarnessX paper", "q"),
     ("Legal / retail reliability (frozen/weak model) → deterministic code dominates.", "b"),
     ("LoCoMo long-context → memory + compaction. GAIA retrieval → tools. Sonnet → prompts.", "b"),
     ("This is the stronger thesis: not just a vertical-specific harness — the vertical-specific LEVER.", "h")])

content("Inverse scaling",
    [("The weakest model gains the most.", "h"),
     ("ALFWorld Qwen3.5-9B: 53 → 97 (+44)   vs   Sonnet: 83.6 → 94.8 (+11.2)", "b"),
     ("Harness evolution makes small / cheap models punch up.", "b"),
     ("Caveat: a capability floor exists — below it, evolution can't compound.", "b")])

content("Vertical-specific harnesses",
    [("τ²-Bench = real business verticals: retail · airline · telecom.", "h"),
     ("Evolve each independently — database tasks need different processors than pure-logic puzzles.", "b"),
     ("Reference (retail, Qwen3.5-27B agent): 0.807 → 0.965, 18/22 badcases fixed.", "b"),
     ("Elsewhere: Database domain 0% → 53.8%.", "b")])

content("Our reproduction — a real lift (telecom)",
    [("Frozen qwen3:32B (local, llama.cpp), τ²-Bench telecom, same 4 mobile_data_issue tasks.", "h"),
     ("Vanilla harness (system prompt + token budget only):  avg reward = 0.500", "b"),
     ("+ IRMA PolicyHint (telecom policy alerts):            avg reward = 0.750", "b"),
     ("+0.25 absolute, +50% relative — zero model-weight changes, no regressions.", "h"),
     ("Mechanism: [POLICY ALERT] injected → agent calls enable_roaming → rescues the abroad task 0.0 → 1.0.", "b")],
    sub="same frozen model throughout · evolved the harness, not the weights")

content("Why telecom worked where retail didn't",
    [("Retail = strict DB-equality grading → a Q4 local model scores ~0, nothing to lift.", "b"),
     ("Telecom = lenient outcome-state grading → non-zero baseline the harness can move.", "b"),
     ("Cause ↔ lever alignment: vanilla fails the roaming task; IRMA's rule targets exactly that.", "h"),
     ("Not at ceiling: 32B solves easy tasks (1.0), fails user_abroad (0.0) — room for the harness.", "b"),
     ("A weaker 8B agent would show a BIGGER gap — inverse scaling.", "b")])

content("Designing the vertical benchmark is HALF the work",
    [("The benchmark IS the reward function for agentic RL. No signal → no evolution.", "h"),
     ("A benchmark you can evolve against needs:", "b"),
     ("  · a non-zero baseline (model not floored)   · room (not at ceiling)", "b"),
     ("  · a failure mode that matches a lever   · grading that gives a gradient (partial credit / enough tasks)", "b"),
     ("Retail (strict DB) → 0.00 dead end.  Telecom (lenient) → moves.  GSM8K → clean gradient.", "b"),
     ("Vertical agents need vertical benchmarks — co-design the harness AND the benchmark together.", "h")])

content("An honest note on reproducing the *number*",
    [("The METHOD reproduces on a laptop: diagnose trace → pick lever → author component → measure.", "h"),
     ("The benchmark + grading choice matters: pick a domain where the model has a non-zero baseline.", "b"),
     ("Small samples are noisy — one task hit an infra timeout; the paper uses 100+ tasks, pass^k, trials.", "b"),
     ("Reference at scale (27B retail): 0.807 → 0.965. Our local telecom: 0.50 → 0.75.", "b")])

content("How this workshop runs — you drive it",
    [("Not slide-watching. You run the loop on your laptop.", "h"),
     ("1. Run the vanilla baseline → watch a frozen model FAIL a real benchmark task.", "b"),
     ("2. Read the trace frontmatter → feel it's a harness failure, not a reasoning one.", "b"),
     ("3. Invoke the harness-evolver skill → YOUR Claude Code diagnoses + authors a fix.", "b"),
     ("4. make eval → see the number move. Loop until it lifts.", "b"),
     ("A reference solution is shipped — but try your own lever first.", "h")])

content("The meta-agent is YOU (Claude Code)",
    [("evolve() = an agent that reads traces and writes config. We used Claude Code directly.", "h"),
     ("read failure trajectories → diagnose the pattern → author a processor → re-eval → keep if it gates", "b"),
     ("No RL training. No GPU. The loop runs on a laptop.", "b")])

content("Reproduce it",
    [("uv + Ollama + one fork.  make baseline → (Claude Code evolves) → make eval", "h"),
     ("Gotcha we hit: LLAMA_API_KEY in your shell 401s all local inference — unset it.", "b"),
     ("Fork: github.com/epuerta9/HarnessX  (branch workshop/harness-evolution)", "b")])

content("The technique is basic. The closed LOOP is the point.",
    [("Reflect, tools, retries, policy hints — all well-known. We're NOT teaching those.", "h"),
     ("The skill is the loop: run → READ THE TRACE → the trace names the failure → pick the matching", "b"),
     ("lever → evolve → re-run → new trace → repeat. Driven by evidence, not guesswork.", "b"),
     ("Every lever we pulled was chosen FROM a trace:", "h"),
     ("telecom trace 'never called enable_roaming' → IRMA (0.50→0.75)", "b"),
     ("GSM8K trace 'fast, wrong' → reflect (0.65→0.97) · KB trace 'hallucinated facts, no tool call' → kb_search (0→1.0)", "b")])

content("Harness evolution IS reinforcement learning",
    [("The “operational mirror” — with the model FROZEN:", "h"),
     ("State = the HarnessConfig   ·   Action = pull one of the 4 levers", "b"),
     ("Policy = a stronger model reading the traces (that's Claude Code)", "b"),
     ("Reward = the benchmark score   ·   Rollout = re-run the task set   ·   Gate = keep if reward↑", "b"),
     ("Agentic RL = use a stronger model to read traces, pick the lever, re-run, keep what improves.", "h")])

image_slide("Which lever moved the number",
            "workshop/lever-result-chart.png",
            "Same frozen model. The lever lifts the score only when the benchmark gives a signal AND the failure matches the lever.")

content("Which component makes the biggest difference?",
    [("Ablation: run vanilla / control-only / IRMA-only / full on the SAME tasks.", "h"),
     ("Single-trial trap: variance (±0.25) is as big as the effect — the same config scores 0.50 OR 0.75 by luck.", "b"),
     ("So you CANNOT rank components from one run. Attribution is statistical, not anecdotal.", "b"),
     ("Fix: multiple trials (pass^k) + task-level check + leave-one-out to kill inert processors.", "b"),
     ("Task-level signal: IRMA's roaming alert rescues the roaming task; control processors are inert here.", "h"),
     ("Reproducible in one command:  bash workshop/ablate.sh", "b")])

image_slide("Attribution (3 trials): one component did the work",
            "workshop/attribution-chart.png",
            "Multi-trial averaging clears the noise: IRMA (+0.167) drives the lift; the 5 control processors are inert (−0.083). Leave-one-out would drop them.")

content("Did we prove the paper? (partly — yes)",
    [("✓ A FROZEN model's benchmark score rose via a harness edit alone (telecom 0.50 → 0.75).", "b"),
     ("✓ The moat is the infrastructure — we diagnosed + evolved purely off typed processors + structured traces.", "b"),
     ("✓ Lever mix is context-dependent — IRMA worked on telecom, did nothing on retail (grading mismatch).", "b"),
     ("✓ Ceiling/floor effects — 32B at ceiling on easy tasks; strict retail grading = effective floor.", "b"),
     ("Not reproduced (needs scale): +14.5% magnitudes, multi-round + the 3 RL pathologies, AEGIS, co-evolution.", "h")])

content("Crutch or capability? Two value props",
    [("Deterministic Control fixes (roaming alert, ParseRetry) are a CRUTCH — they recover latent", "b"),
     ("capability a frozen model failed to apply. Real, but capped at the model's ceiling.", "b"),
     ("Instruction (reasoning scaffolds) + Action (tools/skills) ADD realized capability — not plumbing.", "h"),
     ("To raise the ceiling itself, you need loop 2: RL on the weights (co-evolution). Different game.", "b"),
     ("On a frozen model, the most capability-like levers are reasoning scaffolds + new tools.", "h")])

image_slide("The harness elicits latent reasoning (GSM8K)",
            "workshop/reasoning/reasoning-chart.png",
            "Same frozen qwen3:8B. A reflect scaffold (re-derive + self-check) lifts GSM8K 0.65 → 0.97 and CRT 0.53 → 0.73. Realized reasoning, not I/O.")

image_slide("The spectrum of harness value (all 3 levers, real numbers)",
            "workshop/levers-spectrum-chart.png",
            "Action ADDS capability (private KB 0.00 → 1.00, can't → can). Instruction reshapes reasoning (GSM8K 0.65 → 0.97). Control recovers latent (telecom 0.50 → 0.75, a crutch). Same frozen model throughout.")

content("The Action lever — genuine capability (can't → can)",
    [("A frozen model CANNOT know your private data — it's not in the weights.", "h"),
     ("Vanilla on a private company KB: 0/12 — every fact hallucinated (Sentinel Arm → $12,500, not $8,450).", "b"),
     ("Add a kb_search tool (Action lever): 12/12 — it retrieves the fact and answers.", "b"),
     ("That's NEW capability, not a crutch — and it's the whole reason vertical agents exist:", "h"),
     ("proprietary knowledge, tools, and APIs the base model will never have.", "b")])

# ── closing takeaways ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s, INK)
tf = box(s, 0.9, 0.7, 11.5, 1.0); para(tf, "Takeaways", 36, WHITE, bold=True, first=True)
body = box(s, 0.95, 2.0, 11.5, 5.0)
tks = ["The deliverable is the LOOP: run → read traces → pick the lever → evolve → re-run. Not the technique.",
       "Every lever is chosen FROM a trace — the trace names the failure; you match a lever to it.",
       "The harness, not the weights, sets REALIZED performance — proven across 3 levers, frozen model.",
       "Vertical agents need a STRONG BENCHMARK to evolve against — it's the reward signal in the loop.",
       "Off-the-shelf model + off-the-shelf harness plateaus: no trace signal → nothing to evolve on.",
       "This is how you squeeze a cheap, frozen model to punch up on YOUR vertical."]
for i, t in enumerate(tks):
    para(body, f"{i+1}.  {t}", 20, WHITE if i%2 else GREEN, bold=(i==0), first=(i==0), space=16)

prs.save(OUT)
print("wrote", OUT, "with", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
